#!/usr/bin/env python3
"""
Script para converter arquivo Markdown em apresentação de slides usando python-pptx
Cores personalizadas: preto e #0097a7 (azul-esverdeado)
"""

import os
import re
import fire
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_color):
    """Converte cor hexadecimal para RGB."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def extract_sections_from_md(md_file):
    """
    Extrai seções do arquivo Markdown baseado em cabeçalhos.
    Retorna uma lista de dicionários com título e conteúdo.
    """
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Dividir o conteúdo por cabeçalhos de nível 1 e 2
    sections = []
    
    # Encontrar o título principal (primeiro cabeçalho H1)
    title_match = re.search(r'^# (.*?)$', content, re.MULTILINE)
    main_title = title_match.group(1) if title_match else "Apresentação"
    
    # Encontrar o subtítulo (segundo cabeçalho H1 ou primeiro H2)
    subtitle = ""
    if len(content.split('# ')) > 1:
        subtitle_match = re.search(r'^# (.*?)$', content.split('# ')[1:][0], re.MULTILINE)
        if subtitle_match:
            subtitle = subtitle_match.group(1)
    
    if not subtitle:
        subtitle_match = re.search(r'^## (.*?)$', content, re.MULTILINE)
        if subtitle_match:
            subtitle = subtitle_match.group(1)
    
    # Adicionar slide de título
    sections.append({
        'type': 'title',
        'title': main_title,
        'subtitle': subtitle
    })
    
    # Adicionar slide de agenda
    sections.append({
        'type': 'agenda',
        'title': 'Agenda',
        'content': extract_agenda(content)
    })
    
    # Extrair seções baseadas em cabeçalhos H2
    h2_pattern = r'^## (.*?)$'
    h2_titles = re.findall(h2_pattern, content, re.MULTILINE)
    
    h2_contents = re.split(h2_pattern, content, flags=re.MULTILINE)[1:]  # Primeiro item é o conteúdo antes do primeiro H2
    
    # Combinar títulos e conteúdos
    for i in range(0, len(h2_titles)):
        title = h2_titles[i]
        
        # Obter conteúdo até o próximo H2
        if i < len(h2_contents):
            section_content = h2_contents[i]
            
            # Extrair subseções (H3)
            subsections = extract_subsections(section_content)
            
            if subsections:
                # Se houver subseções, criar um slide para cada uma
                for subsection in subsections:
                    sections.append({
                        'type': 'content',
                        'title': title + " - " + subsection['title'],
                        'content': subsection['content']
                    })
            else:
                # Se não houver subseções, criar um slide para a seção inteira
                sections.append({
                    'type': 'content',
                    'title': title,
                    'content': section_content
                })
    
    return sections

def extract_agenda(content):
    """
    Extrai itens de agenda baseados em cabeçalhos H2.
    """
    h2_pattern = r'^## (.*?)$'
    h2_titles = re.findall(h2_pattern, content, re.MULTILINE)
    
    agenda_items = []
    for title in h2_titles:
        if title != "Introdução" and title != "Conclusão" and "Pré-requisitos" not in title:
            agenda_items.append(title)
    
    return agenda_items

def extract_subsections(content):
    """
    Extrai subseções baseadas em cabeçalhos H3.
    """
    h3_pattern = r'^### (.*?)$'
    h3_titles = re.findall(h3_pattern, content, re.MULTILINE)
    
    if not h3_titles:
        return []
    
    h3_contents = re.split(h3_pattern, content, flags=re.MULTILINE)[1:]  # Primeiro item é o conteúdo antes do primeiro H3
    
    subsections = []
    for i in range(0, len(h3_titles)):
        title = h3_titles[i]
        
        # Obter conteúdo até o próximo H3
        if i < len(h3_contents):
            subsection_content = h3_contents[i]
            subsections.append({
                'title': title,
                'content': subsection_content
            })
    
    return subsections

def extract_bullet_points(content):
    """
    Extrai pontos de marcadores do conteúdo.
    """
    # Remover blocos de código
    content = re.sub(r'```.*?```', '', content, flags=re.DOTALL)
    
    # Extrair listas não ordenadas
    bullet_points = []
    
    # Padrão para listas com marcadores
    list_pattern = r'^\s*[\*\-\+]\s+(.*?)$'
    bullet_matches = re.findall(list_pattern, content, re.MULTILINE)
    
    if bullet_matches:
        bullet_points.extend(bullet_matches)
    
    # Se não houver marcadores explícitos, extrair parágrafos curtos
    if not bullet_points:
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        bullet_points = [p for p in paragraphs if len(p) < 200 and not p.startswith('#')]
    
    return bullet_points

def extract_code_blocks(content):
    """
    Extrai blocos de código do conteúdo.
    """
    code_pattern = r'```(?:python)?\n(.*?)\n```'
    code_blocks = re.findall(code_pattern, content, flags=re.DOTALL)
    return code_blocks

def extract_quotes(content):
    """
    Extrai citações do conteúdo.
    """
    quote_pattern = r'^>\s*(.*?)$'
    quotes = re.findall(quote_pattern, content, re.MULTILINE)
    return quotes

def create_presentation(sections, output_file, main_color="#0097a7"):
    """
    Cria uma apresentação PowerPoint a partir das seções extraídas.
    
    Args:
        sections: Lista de dicionários com informações das seções
        output_file: Caminho para o arquivo de saída
        main_color: Cor principal em formato hexadecimal
    
    Returns:
        Caminho para o arquivo de apresentação gerado
    """
    prs = Presentation()
    
    # Converter cor hexadecimal para RGB
    main_color_rgb = hex_to_rgb(main_color)
    black = RGBColor(0, 0, 0)
    white = RGBColor(255, 255, 255)
    gray = RGBColor(128, 128, 128)
    light_gray = RGBColor(240, 240, 240)
    
    # Definir tamanho dos slides (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Criar slide de título
    title_slide_layout = prs.slide_layouts[0]  # Layout de título
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Configurar título
    title = slide.shapes.title
    title.text = sections[0]['title']
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = main_color_rgb
    title.text_frame.paragraphs[0].font.bold = True
    
    # Configurar subtítulo
    subtitle = slide.placeholders[1]
    subtitle.text = sections[0]['subtitle']
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle.text_frame.paragraphs[0].font.color.rgb = black
    
    # Adicionar rodapé com informações
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5))
    footer_text = footer.text_frame
    p = footer_text.paragraphs[0]
    p.text = "Disciplina de MLOps"
    p.font.size = Pt(14)
    p.font.color.rgb = gray
    p.alignment = PP_ALIGN.RIGHT
    
    # Processar as demais seções
    for section in sections[1:]:
        if section['type'] == 'agenda':
            # Criar slide de agenda
            content_slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Configurar título
            title = slide.shapes.title
            title.text = section['title']
            title.text_frame.paragraphs[0].font.size = Pt(40)
            title.text_frame.paragraphs[0].font.color.rgb = main_color_rgb
            title.text_frame.paragraphs[0].font.bold = True
            
            # Adicionar itens de agenda
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            for item in section['content']:
                p = tf.add_paragraph()
                p.text = "• " + item
                p.font.size = Pt(28)
                p.font.color.rgb = black
                p.space_after = Pt(12)
        
        elif section['type'] == 'content':
            # Criar slide de conteúdo
            content_slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Configurar título
            title = slide.shapes.title
            title.text = section['title']
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.color.rgb = main_color_rgb
            title.text_frame.paragraphs[0].font.bold = True
            
            # Extrair pontos de marcadores
            bullet_points = extract_bullet_points(section['content'])
            
            # Extrair citações
            quotes = extract_quotes(section['content'])
            
            # Extrair blocos de código
            code_blocks = extract_code_blocks(section['content'])
            
            # Adicionar conteúdo
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Adicionar citação se houver
            if quotes:
                p = tf.add_paragraph()
                p.text = '"' + quotes[0] + '"'
                p.font.italic = True
                p.font.size = Pt(20)
                p.font.color.rgb = gray
                p.space_after = Pt(20)
            
            # Adicionar pontos de marcadores
            for point in bullet_points[:5]:  # Limitar a 5 pontos por slide
                p = tf.add_paragraph()
                p.text = "• " + point
                p.font.size = Pt(24)
                p.font.color.rgb = black
                p.space_after = Pt(12)
            
            # Adicionar bloco de código se houver
            if code_blocks:
                # Adicionar caixa de texto para o código
                left = Inches(1.0)
                top = Inches(4.0)
                width = Inches(11.33)
                height = Inches(2.5)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf_code = textbox.text_frame
                
                # Limitar o código a algumas linhas para caber no slide
                code_lines = code_blocks[0].split('\n')[:8]  # Máximo de 8 linhas
                code_text = '\n'.join(code_lines)
                
                p = tf_code.add_paragraph()
                p.text = code_text
                p.font.name = 'Courier New'
                p.font.size = Pt(16)
                
                # Adicionar fundo cinza claro para o código
                fill = textbox.fill
                fill.solid()
                fill.fore_color.rgb = light_gray
    
    # Adicionar slide final
    final_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(final_slide_layout)
    
    # Configurar título
    title = slide.shapes.title
    title.text = "Perguntas e Recursos Adicionais"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = main_color_rgb
    title.text_frame.paragraphs[0].font.bold = True
    
    # Adicionar conteúdo
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "• Recapitulação dos conceitos principais"
    p.font.size = Pt(28)
    p.font.color.rgb = black
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Próximos passos: demonstração prática com scikit-learn e Airflow"
    p.font.size = Pt(28)
    p.font.color.rgb = black
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Recursos para aprofundamento"
    p.font.size = Pt(28)
    p.font.color.rgb = black
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Contato para dúvidas"
    p.font.size = Pt(28)
    p.font.color.rgb = black
    p.space_after = Pt(12)
    
    # Salvar apresentação
    prs.save(output_file)
    print(f"Apresentação criada com sucesso: {output_file}")
    return output_file

def create_slides_from_structure(structure_file, output_file, main_color="#0097a7"):
    """
    Cria uma apresentação PowerPoint a partir de um arquivo de estrutura de slides.
    
    Args:
        structure_file: Caminho para o arquivo de estrutura de slides
        output_file: Caminho para o arquivo de saída
        main_color: Cor principal em formato hexadecimal
    
    Returns:
        Caminho para o arquivo de apresentação gerado
    """
    with open(structure_file, 'r', encoding='utf-8') as f:
        structure_content = f.read()
    
    # Extrair seções de slides
    slide_sections = []
    
    # Padrão para encontrar definições de slides
    slide_pattern = r'^## Slide (\d+): (.*?)$(.*?)(?=^## Slide \d+:|$)'
    slide_matches = re.findall(slide_pattern, structure_content, re.MULTILINE | re.DOTALL)
    
    # Slide de título
    title_match = re.search(r'## Slide 1: (.*?)$(.*?)(?=^## Slide \d+:|$)', structure_content, re.MULTILINE | re.DOTALL)
    if title_match:
        title = title_match.group(1).strip()
        content = title_match.group(2).strip()
        
        # Extrair subtítulo
        subtitle = ""
        subtitle_match = re.search(r'- Subtítulo: (.*?)$', content, re.MULTILINE)
        if subtitle_match:
            subtitle = subtitle_match.group(1).strip()
        
        slide_sections.append({
            'type': 'title',
            'title': title,
            'subtitle': subtitle
        })
    
    # Slide de agenda
    agenda_match = re.search(r'## Slide 2: (.*?)$(.*?)(?=^## Slide \d+:|$)', structure_content, re.MULTILINE | re.DOTALL)
    if agenda_match:
        title = agenda_match.group(1).strip()
        content = agenda_match.group(2).strip()
        
        # Extrair itens de agenda
        agenda_items = []
        items_match = re.findall(r'- (.*?)$', content, re.MULTILINE)
        for item in items_match:
            if "Tópicos" not in item and "Imagem" not in item:
                agenda_items.append(item)
        
        slide_sections.append({
            'type': 'agenda',
            'title': title,
            'content': agenda_items
        })
    
    # Demais slides
    for slide_num, title, content in slide_matches[2:]:  # Pular os dois primeiros slides (título e agenda)
        slide_sections.append({
            'type': 'content',
            'title': title.strip(),
            'content': content.strip()
        })
    
    # Criar apresentação
    return create_presentation(slide_sections, output_file, main_color)

def md_to_slides(md_file, output_file=None, main_color="#0097a7"):
    """
    Converte um arquivo Markdown em uma apresentação de slides.
    
    Args:
        md_file: Caminho para o arquivo Markdown de entrada
        output_file: Caminho para o arquivo PowerPoint de saída (opcional)
        main_color: Cor principal em formato hexadecimal (padrão: #0097a7)
    
    Returns:
        Caminho para o arquivo de apresentação gerado
    """
    # Verificar se o arquivo de entrada existe
    if not os.path.exists(md_file):
        raise FileNotFoundError(f"Arquivo {md_file} não encontrado.")
    
    # Definir nome do arquivo de saída se não fornecido
    if output_file is None:
        base_name = os.path.splitext(os.path.basename(md_file))[0]
        output_file = f"{base_name}_Slides.pptx"
    
    # Extrair seções do arquivo Markdown
    print(f"Extraindo seções de: {md_file}")
    sections = extract_sections_from_md(md_file)
    
    # Criar apresentação
    print(f"Criando apresentação: {output_file}")
    return create_presentation(sections, output_file, main_color)

def structure_to_slides(structure_file, output_file=None, main_color="#0097a7"):
    """
    Cria uma apresentação de slides a partir de um arquivo de estrutura.
    
    Args:
        structure_file: Caminho para o arquivo de estrutura de slides
        output_file: Caminho para o arquivo PowerPoint de saída (opcional)
        main_color: Cor principal em formato hexadecimal (padrão: #0097a7)
    
    Returns:
        Caminho para o arquivo de apresentação gerado
    """
    # Verificar se o arquivo de entrada existe
    if not os.path.exists(structure_file):
        raise FileNotFoundError(f"Arquivo {structure_file} não encontrado.")
    
    # Definir nome do arquivo de saída se não fornecido
    if output_file is None:
        base_name = os.path.splitext(os.path.basename(structure_file))[0]
        output_file = f"{base_name}_Slides.pptx"
    
    # Criar apresentação a partir da estrutura
    print(f"Criando slides a partir da estrutura: {structure_file}")
    return create_slides_from_structure(structure_file, output_file, main_color)

if __name__ == "__main__":
    fire.Fire({
        'md': md_to_slides,
        'structure': structure_to_slides
    })
